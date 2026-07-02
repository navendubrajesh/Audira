"""Document ingestion — txt/md/docx/pdf/pptx (TCA-007)."""

from __future__ import annotations

import io
from dataclasses import dataclass


@dataclass
class ParsedDocument:
    text: str
    format: str
    page_count: int
    slide_count: int
    warnings: list[str]


def parse_document(filename: str, content: bytes, content_type: str | None = None) -> ParsedDocument:
    name = (filename or "").lower()
    warnings: list[str] = []

    if name.endswith((".txt", ".md")) or (content_type or "").startswith("text/"):
        text = content.decode("utf-8", errors="ignore")
        return ParsedDocument(text=text, format="text", page_count=1, slide_count=0, warnings=warnings)

    if name.endswith(".docx") or "wordprocessingml" in (content_type or ""):
        return _parse_docx(content, warnings)

    if name.endswith(".pdf") or content_type == "application/pdf":
        return _parse_pdf(content, warnings)

    if name.endswith((".pptx", ".ppt")) or "presentationml" in (content_type or ""):
        return _parse_pptx(content, warnings)

    text = content.decode("utf-8", errors="ignore")
    if not text.strip():
        warnings.append("Unknown format — could not extract text.")
    return ParsedDocument(text=text, format="unknown", page_count=0, slide_count=0, warnings=warnings)


def _parse_docx(content: bytes, warnings: list[str]) -> ParsedDocument:
    try:
        from docx import Document

        doc = Document(io.BytesIO(content))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        return ParsedDocument(
            text="\n\n".join(parts),
            format="docx",
            page_count=max(1, len(parts) // 20),
            slide_count=0,
            warnings=warnings,
        )
    except ImportError:
        warnings.append("python-docx not installed — install for DOCX support.")
        return ParsedDocument(text="", format="docx", page_count=0, slide_count=0, warnings=warnings)
    except Exception as exc:
        warnings.append(f"DOCX parse error: {exc}")
        return ParsedDocument(text="", format="docx", page_count=0, slide_count=0, warnings=warnings)


def _parse_pdf(content: bytes, warnings: list[str]) -> ParsedDocument:
    try:
        import fitz

        doc = fitz.open(stream=content, filetype="pdf")
        pages = [page.get_text() for page in doc]
        return ParsedDocument(
            text="\n\n".join(pages),
            format="pdf",
            page_count=len(pages),
            slide_count=0,
            warnings=warnings,
        )
    except ImportError:
        warnings.append("PyMuPDF not installed — install for PDF support.")
        return ParsedDocument(text="", format="pdf", page_count=0, slide_count=0, warnings=warnings)
    except Exception as exc:
        warnings.append(f"PDF parse error: {exc}")
        return ParsedDocument(text="", format="pdf", page_count=0, slide_count=0, warnings=warnings)


def _parse_pptx(content: bytes, warnings: list[str]) -> ParsedDocument:
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        slides: list[str] = []
        for slide in prs.slides:
            lines: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines.append(shape.text)
            if lines:
                slides.append("\n".join(lines))
        return ParsedDocument(
            text="\n\n---\n\n".join(slides),
            format="pptx",
            page_count=0,
            slide_count=len(slides),
            warnings=warnings,
        )
    except ImportError:
        warnings.append("python-pptx not installed — install for PPTX support.")
        return ParsedDocument(text="", format="pptx", page_count=0, slide_count=0, warnings=warnings)
    except Exception as exc:
        warnings.append(f"PPTX parse error: {exc}")
        return ParsedDocument(text="", format="pptx", page_count=0, slide_count=0, warnings=warnings)
